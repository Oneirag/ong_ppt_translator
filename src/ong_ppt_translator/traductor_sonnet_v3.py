from ong_ppt_translator.translate_text import translate_text_with_openai
import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR_INDEX
from ong_ppt_translator.process_runs import parse_html_text
from ong_ppt_translator import logger


def limpiar_etiquetas_vacias(texto):
    # Patrón para encontrar etiquetas que contienen solo espacios o etiquetas vacías
    patron = r'<([a-z][a-z0-9]*)\b[^>]*>((\s*)|(<([a-z][a-z0-9]*)\b[^>]*>(\s*)</\5>))</\1>'

    # Aplicar reemplazo iterativamente hasta que no haya más cambios
    texto_anterior = None
    texto_actual = texto

    while texto_anterior != texto_actual:
        texto_anterior = texto_actual
        texto_actual = re.sub(patron, r'\2', texto_anterior)

    return texto_actual

def limpiar_html(texto):
    """Removes all HTML tags but b, i and u"""
    retval = re.sub(r"</?(?!b|i|u\b)[a-zA-Z0-9]+.*?>", "", texto)
    retval = limpiar_etiquetas_vacias(retval)
    return retval


def extract_markdown_from_shape(shape):
    """
    Convierte el contenido de un shape a markdown y también devuelve el color de los runs
    """
    if not hasattr(shape, 'text_frame'):
        return [], []

    markdown_paragraphs = []
    color_paragraphs = []

    for paragraph in shape.text_frame.paragraphs:
        # Construir markdown para el párrafo
        paragraph_markdown = ""

        # # Manejar nivel de bullet/indentación
        # if paragraph.level > 0:
        #     paragraph_markdown += "  " * paragraph.level + "- "

        color_paragraphs.append([])
        # Procesar runs
        for run in paragraph.runs:
            # Aplicar formatos
            text = run.text

            # Orden de aplicación de formato importante
            if run.font.bold:
                # text = f"**{text}**"
                text = f"<b>{text}</b>"
            if run.font.italic:
                # text = f"*{text}*"
                text = f"<i>{text}</i>"
            if run.font.underline:
                text = f"<u>{text}</u>"

            paragraph_markdown += text
            color = run.font.color if run.font.color.type else None
            color_paragraphs[-1].append(color)
        markdown_paragraphs.append(paragraph_markdown)

        # # Solo añadir si hay contenido
        # if paragraph_markdown.strip():
        #     markdown_paragraphs.append(paragraph_markdown)

    return markdown_paragraphs, color_paragraphs


def is_text_shape(shape):
    """
    Verifica si un shape puede contener texto
    """
    text_shape_types = [
        MSO_SHAPE_TYPE.TEXT_BOX,
        #MSO_SHAPE_TYPE.RECTANGLE,
        #MSO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        MSO_SHAPE_TYPE.PLACEHOLDER
    ]

    return (
            hasattr(shape, 'text_frame') or
            (hasattr(shape, 'shape_type') and shape.shape_type in text_shape_types)
    )

def iter_shapes(shapes: list):
    for shape in shapes:
        if not shape.shape_type is MSO_SHAPE_TYPE.GROUP:
            yield shape
        else:
            yield from iter_shapes(shape.shapes)



def translate_powerpoint(input_file, output_file, end: int=None, start: int=None):
    """
    Traduce una presentación de PowerPoint usando markdown
    """

    # Cargar presentación
    prs = Presentation(input_file)

    prompt_tokens = 0
    completion_tokens = 0
    llm_time = 0

    # Iterar sobre todas las diapositivas
    for idx_slide, slide in enumerate(prs.slides):
        if (end and idx_slide > end) or (start and idx_slide < start) :
            continue
        logger.info(f"Processing slide {idx_slide + 1}")
        for shape in iter_shapes(slide.shapes):
            # Verificar si es un shape con texto
            if not is_text_shape(shape):
                continue
            try:
                logger.debug(f"Processing {shape.text}")
            except:
                logger.debug("Shape has no text")
            # Extraer markdown
            markdown_paragraphs, color_paragraphs = extract_markdown_from_shape(shape)

            # Saltar si no hay contenido
            if not markdown_paragraphs:
                continue

            # # Limpiar shape original
            # if hasattr(shape, 'text_frame'):
            #     shape.text_frame.clear()

            # Traducir y reformatear
            for idx, markdown_paragraph in enumerate(markdown_paragraphs):
                # Traducir texto
                translation = translate_text_with_openai(markdown_paragraph)
                prompt_tokens += translation.prompt_tokens
                completion_tokens += translation.completion_tokens
                llm_time += translation.seconds

                # Convertir markdown traducido a runs
                translated_runs = parse_html_text(translation.text)

                # Añadir párrafo traducido
                if hasattr(shape, 'text_frame'):
                    if idx + 1 > len(shape.text_frame.paragraphs):
                        new_paragraph = shape.text_frame.add_paragraph()
                    else:
                        new_paragraph = shape.text_frame.paragraphs[idx]

                    # Borrar el texto de los run actuales
                    for run in new_paragraph.runs:
                        run.text = ""

                    # Añadir runs
                    for idx_run, run_info in enumerate(translated_runs):
                        if idx_run < len(new_paragraph.runs):
                            new_run = new_paragraph.runs[idx_run]
                        else:
                            new_run = new_paragraph.add_run()
                        # new_run = new_paragraph.add_run()
                        new_run.text = run_info['text']

                        # Aplicar formato
                        new_run.font.bold = run_info['bold']
                        new_run.font.italic = run_info['italic']
                        new_run.font.underline = run_info['underline']

                        try:
                            run_color = color_paragraphs[idx][idx_run]
                        except Exception as e:
                            run_color = None
                        if run_color:
                            match run_color.type:
                                case MSO_COLOR_TYPE.RGB:
                                    attr = "rgb"
                                case MSO_COLOR_TYPE.PRESET:
                                    if run_color.theme_color == MSO_THEME_COLOR_INDEX.NOT_THEME_COLOR:
                                        attr = None
                                    else:
                                        attr = "theme_color"
                                case _:
                                    attr = None
                            if attr:
                                try:
                                    setattr(new_run.font.color, attr, getattr(run_color, attr))
                                except Exception as e:
                                    logger.exception(e)
                            # for attr in "theme_color", "rgb":
                            if new_run.font.color.type != run_color.type:
                                logger.warning(f"Error changing color of {new_run.text}")

    # Guardar presentación traducida
    prs.save(output_file)
    logger.info(f"Translation took {prompt_tokens=} {completion_tokens=} in {llm_time:.2f} seconds")


if __name__ == '__main__':

    pass
    # Uso del script
    #input_file = 'presentacion_original.pptx'
    # output_file = 'presentacion_traducida.pptx'
    # translate_powerpoint(input_file, output_file, start=22-1, end=22-1) #, end=7) #, start=6)

    #print(f"Presentación traducida guardada en {output_file}")